#!/usr/bin/env python3
"""Build an editable presentation for the current joint model architecture."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUTPUT = ROOT / "docs/presentations/dlcpd25-joint-model-architecture.pptx"

W = Inches(13.333)
H = Inches(7.5)
FONT = "Noto Sans CJK SC"

INK = RGBColor(31, 41, 55)
MUTED = RGBColor(91, 102, 116)
PAPER = RGBColor(247, 249, 250)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(34, 113, 85)
GREEN_LIGHT = RGBColor(220, 239, 230)
BLUE = RGBColor(38, 92, 143)
BLUE_LIGHT = RGBColor(222, 235, 247)
RED = RGBColor(181, 65, 55)
RED_LIGHT = RGBColor(249, 226, 222)
GOLD = RGBColor(184, 127, 35)
GOLD_LIGHT = RGBColor(249, 238, 208)
GRAY_LIGHT = RGBColor(231, 235, 238)


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    size: float = 18,
    color: RGBColor = INK,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.MIDDLE,
):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = Inches(0.04)
    frame.margin_top = frame.margin_bottom = Inches(0.02)
    frame.vertical_anchor = valign
    paragraph = frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_box(
    slide,
    title: str,
    subtitle: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: RGBColor,
    line: RGBColor,
    title_size: float = 18,
    subtitle_size: float = 11.5,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1.5)
    shape.adjustments[0] = 0.08
    add_text(
        slide,
        title,
        x + 0.12,
        y + 0.09,
        w - 0.24,
        0.35,
        size=title_size,
        color=line,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    if subtitle:
        add_text(
            slide,
            subtitle,
            x + 0.14,
            y + 0.48,
            w - 0.28,
            h - 0.56,
            size=subtitle_size,
            color=INK,
            align=PP_ALIGN.CENTER,
        )
    return shape


def add_arrow(slide, x1: float, y1: float, x2: float, y2: float, color=INK, width=2.2):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2)
    )
    line.line.color.rgb = color
    line.line.width = Pt(width)
    line.line.end_arrowhead = True
    return line


def add_title(slide, title: str, subtitle: str = ""):
    add_text(slide, title, 0.55, 0.25, 12.2, 0.5, size=25, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.58, 0.78, 12.0, 0.3, size=11.5, color=MUTED)
    rule = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.58), Inches(1.14), Inches(12.15), Inches(0.03)
    )
    rule.fill.solid()
    rule.fill.fore_color.rgb = GREEN
    rule.line.fill.background()


def add_footer(slide, page: int):
    add_text(
        slide,
        "DLCPD-25 + IP102 | 当前代码架构",
        0.58,
        7.13,
        4.2,
        0.2,
        size=8.5,
        color=MUTED,
    )
    add_text(
        slide,
        str(page),
        12.2,
        7.1,
        0.5,
        0.22,
        size=8.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def set_background(slide, color=PAPER):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_bullet_list(slide, items: list[str], x: float, y: float, w: float, h: float, size=15):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = item
        paragraph.level = 0
        paragraph.font.name = FONT
        paragraph.font.size = Pt(size)
        paragraph.font.color.rgb = INK
        paragraph.space_after = Pt(8)
        paragraph.text = "• " + item
    return box


def build() -> None:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    blank = prs.slide_layouts[6]

    # Slide 1: title
    slide = prs.slides.add_slide(blank)
    set_background(slide, WHITE)
    block = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(4.7), H
    )
    block.fill.solid()
    block.fill.fore_color.rgb = GREEN
    block.line.fill.background()
    add_text(slide, "MODEL\nARCHITECTURE", 0.58, 0.65, 3.55, 1.5, size=30, color=WHITE, bold=True)
    add_text(slide, "共享主干 · 双任务头 · 单一权重", 0.62, 2.35, 3.55, 0.55, size=13.5, color=GREEN_LIGHT)
    add_text(slide, "基于DLCPD-25数据集的\n农产品病虫害与缺陷分类目标检测系统", 5.25, 1.0, 7.25, 1.6, size=27, bold=True)
    add_text(slide, "当前实现：ResNet-50 + FPN + Faster R-CNN", 5.3, 2.82, 6.9, 0.5, size=18, color=BLUE, bold=True)
    add_box(slide, "203类整图分类", "DLCPD-25", 5.3, 3.72, 2.25, 1.25, fill=GREEN_LIGHT, line=GREEN)
    add_box(slide, "96类害虫检测", "IP102边界框", 7.85, 3.72, 2.25, 1.25, fill=BLUE_LIGHT, line=BLUE)
    add_box(slide, "统一输入", "RGB 224 × 224", 10.4, 3.72, 2.0, 1.25, fill=GOLD_LIGHT, line=GOLD)
    add_text(slide, "状态：J1/J2 已通过；J3 正式联合训练尚未启动", 5.3, 5.6, 7.0, 0.45, size=14, color=RED, bold=True)
    add_text(slide, "生成日期：2026-08-12", 5.3, 6.3, 4.0, 0.3, size=10.5, color=MUTED)

    # Slide 2: overall architecture
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "1. 总体模型架构", "一张输入、一次共享主干前向、两个任务输出、一个联合 checkpoint")
    add_box(slide, "输入图片", "任意尺寸 RGB", 0.55, 2.65, 1.35, 1.0, fill=GOLD_LIGHT, line=GOLD)
    add_arrow(slide, 1.9, 3.15, 2.25, 3.15, GOLD)
    add_box(slide, "统一预处理", "EXIF校正 → RGB\nBicubic直缩 224×224\nImageNet mean/std", 2.25, 2.35, 2.05, 1.6, fill=WHITE, line=GOLD, subtitle_size=11)
    add_arrow(slide, 4.3, 3.15, 4.65, 3.15, GREEN)
    add_box(slide, "共享 ResNet-50 主干", "C2 / C3 / C4 / C5\n只计算一次", 4.65, 2.25, 2.25, 1.8, fill=GREEN_LIGHT, line=GREEN, title_size=19)
    add_arrow(slide, 6.9, 2.8, 7.45, 2.05, GREEN)
    add_arrow(slide, 6.9, 3.5, 7.45, 4.35, GREEN)
    add_box(slide, "分类分支", "C5(2048×7×7)\n全局平均池化\nLinear 2048→203", 7.45, 1.35, 2.2, 1.7, fill=GREEN_LIGHT, line=GREEN)
    add_box(slide, "检测分支", "FPN → RPN → ROI Heads\nFaster R-CNN\n96前景类 + 背景", 7.45, 3.6, 2.2, 1.8, fill=BLUE_LIGHT, line=BLUE)
    add_arrow(slide, 9.65, 2.2, 10.1, 2.2, GREEN)
    add_arrow(slide, 9.65, 4.5, 10.1, 4.5, BLUE)
    add_box(slide, "分类输出", "宿主 / 四大类 / 具体类别\nTop-5 + 置信度", 10.1, 1.45, 2.65, 1.5, fill=WHITE, line=GREEN)
    add_box(slide, "检测输出", "害虫框 + 分数\n映射为DLCPD-25 class_id", 10.1, 3.75, 2.65, 1.5, fill=WHITE, line=BLUE)
    add_text(slide, "联合权重：共享主干 + 203类分类头 + FPN/RPN/ROI检测头", 3.1, 6.05, 7.2, 0.42, size=15, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 2)

    # Slide 3: internals
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "2. 共享特征与双分支细节", "分类直接使用 C5；检测使用 C2-C5 构建多尺度 FPN")
    stages = [
        ("Stem", "7×7 Conv + MaxPool\n输出 64×56×56", 0.55, GRAY_LIGHT, MUTED),
        ("C2", "layer1\n256×56×56", 2.3, GREEN_LIGHT, GREEN),
        ("C3", "layer2\n512×28×28", 4.05, GREEN_LIGHT, GREEN),
        ("C4", "layer3\n1024×14×14", 5.8, GREEN_LIGHT, GREEN),
        ("C5", "layer4\n2048×7×7", 7.55, GREEN_LIGHT, GREEN),
    ]
    for index, (name, subtitle, x, fill, line) in enumerate(stages):
        add_box(slide, name, subtitle, x, 2.0, 1.4, 1.45, fill=fill, line=line, title_size=17, subtitle_size=10.5)
        if index < len(stages) - 1:
            add_arrow(slide, x + 1.4, 2.72, x + 1.75, 2.72, GREEN)
    add_arrow(slide, 8.95, 2.25, 9.55, 1.62, GREEN)
    add_box(slide, "分类头", "AdaptiveAvgPool\n2048维向量\nFC → 203 logits", 9.55, 1.05, 2.55, 1.6, fill=GREEN_LIGHT, line=GREEN)

    # FPN band
    add_text(slide, "多尺度检测路径", 0.58, 4.05, 1.7, 0.35, size=16, color=BLUE, bold=True)
    fpn_x = [2.3, 4.05, 5.8, 7.55]
    fpn_names = ["P2\n256×56×56", "P3\n256×28×28", "P4\n256×14×14", "P5/P6\n256×7×7 / 4×4"]
    for x, name in zip(fpn_x, fpn_names):
        add_box(slide, name.split("\n")[0], "\n".join(name.split("\n")[1:]), x, 4.45, 1.4, 1.15, fill=BLUE_LIGHT, line=BLUE, title_size=15, subtitle_size=9.5)
    for x in [3.0, 4.75, 6.5, 8.25]:
        add_arrow(slide, x, 3.46, x, 4.43, BLUE, 1.5)
    add_arrow(slide, 8.95, 5.02, 9.35, 5.02, BLUE)
    add_box(slide, "RPN", "候选区域", 9.35, 4.55, 1.15, 0.95, fill=WHITE, line=BLUE, title_size=15, subtitle_size=9.5)
    add_arrow(slide, 10.5, 5.02, 10.82, 5.02, BLUE)
    add_box(slide, "ROI Heads", "框回归 + 97类分类", 10.82, 4.4, 1.95, 1.25, fill=BLUE_LIGHT, line=BLUE, title_size=15, subtitle_size=9.5)
    add_text(slide, "注意：检测内部97类 = 96个害虫前景类 + 1个背景类", 7.3, 6.25, 5.45, 0.35, size=12, color=BLUE, bold=True, align=PP_ALIGN.RIGHT)
    add_footer(slide, 3)

    # Slide 4: training
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "3. 双数据集交替联合训练", "不是把两种标签混成一个 batch，而是分类 step 与检测 step 固定 1:1 交替")
    add_box(slide, "J1 初始化", "DLCPD-25分类权重\nVal Top-1 90.7882%", 0.6, 1.55, 2.0, 1.25, fill=GOLD_LIGHT, line=GOLD)
    add_arrow(slide, 2.6, 2.18, 3.05, 2.18, GOLD)
    add_box(slide, "构建联合模型", "加载共享主干与分类头\n随机初始化检测分支", 3.05, 1.48, 2.55, 1.4, fill=WHITE, line=GREEN)
    add_arrow(slide, 5.6, 2.18, 6.0, 2.18, GREEN)
    add_box(slide, "同一联合 checkpoint", "全部参数 + optimizer\n scheduler + AMP + RNG", 6.0, 1.48, 2.5, 1.4, fill=GREEN_LIGHT, line=GREEN)

    add_box(slide, "DLCPD-25 batch", "分类标签\n177,021张 train", 0.65, 4.05, 2.0, 1.3, fill=GREEN_LIGHT, line=GREEN)
    add_arrow(slide, 2.65, 4.7, 3.15, 4.7, GREEN)
    add_box(slide, "分类 step", "交叉熵\n更新：共享主干 + 分类头\n冻结：检测分支", 3.15, 3.85, 2.55, 1.7, fill=WHITE, line=GREEN)
    add_box(slide, "IP102 batch", "边界框标签\n12,142张 train", 7.65, 4.05, 2.0, 1.3, fill=BLUE_LIGHT, line=BLUE)
    add_arrow(slide, 9.65, 4.7, 10.15, 4.7, BLUE)
    add_box(slide, "检测 step", "RPN + ROI损失\n更新：共享主干 + 检测头\n冻结：分类头", 10.15, 3.85, 2.55, 1.7, fill=WHITE, line=BLUE)
    add_arrow(slide, 5.7, 4.7, 7.65, 4.7, RED, 2.5)
    add_text(slide, "1 : 1 交替", 6.1, 4.23, 1.15, 0.35, size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "分类 AMP FP16", 3.45, 5.82, 1.9, 0.32, size=11.5, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "检测 FP32（避免 NaN）", 10.25, 5.82, 2.3, 0.32, size=11.5, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 4)

    # Slide 5: optimization and selection
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "4. 参数更新、学习率与模型选择", "保护分类能力，同时让随机初始化的检测头更快适配 IP102")
    add_box(slide, "共享主干", "ResNet-50 body\n学习率 1e-5\n分类/检测 step 都更新", 0.7, 1.65, 3.0, 1.65, fill=GREEN_LIGHT, line=GREEN, title_size=20)
    add_box(slide, "203类分类头", "Linear(2048→203)\n学习率 1e-5\n仅分类 step 更新", 4.0, 1.65, 2.7, 1.65, fill=GOLD_LIGHT, line=GOLD, title_size=20)
    add_box(slide, "检测分支", "FPN + RPN + ROI Heads\n学习率 1e-4\n仅检测 step 更新", 7.0, 1.65, 2.8, 1.65, fill=BLUE_LIGHT, line=BLUE, title_size=20)
    add_box(slide, "优化器", "AdamW + Cosine Scheduler\nWeight decay 1e-4\n正式计划 5 epochs", 10.1, 1.65, 2.55, 1.65, fill=WHITE, line=MUTED, title_size=20)
    add_text(slide, "每个 epoch 分别评估两个 val 集", 0.75, 4.0, 4.2, 0.4, size=18, color=INK, bold=True)
    add_bullet_list(
        slide,
        [
            "分类门槛：Top-1 ≥ 88.7837%（相对 J1 最多下降 2 个百分点）",
            "分类安全线：Top-1 < 85% 自动停止，避免灾难性遗忘",
            "门槛合格的 checkpoint 中，选择检测 mAP@0.5:0.95 最高者",
            "J3 不读取 test；J4 冻结后分类/检测 test 各执行一次",
        ],
        0.72,
        4.48,
        11.8,
        1.85,
        size=15,
    )
    add_footer(slide, 5)

    # Slide 6: inference and boundaries
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_title(slide, "5. 最终推理结果与能力边界", "系统是一份联合权重，但两个任务的数据监督范围不同")
    add_box(slide, "上传一张图片", "原图保留用于展示", 0.65, 2.4, 1.7, 1.2, fill=GOLD_LIGHT, line=GOLD)
    add_arrow(slide, 2.35, 3.0, 2.8, 3.0, GOLD)
    add_box(slide, "联合模型", "一个224输入\n一个ResNet-50主干\n一次联合调用", 2.8, 2.1, 2.4, 1.8, fill=GREEN_LIGHT, line=GREEN, title_size=20)
    add_arrow(slide, 5.2, 2.55, 5.8, 1.9, GREEN)
    add_arrow(slide, 5.2, 3.45, 5.8, 4.2, BLUE)
    add_box(slide, "整图分类结果", "覆盖全部203类\n害虫、病害、健康、缺陷\nTop-5 + 置信度", 5.8, 1.25, 2.75, 1.8, fill=GREEN_LIGHT, line=GREEN)
    add_box(slide, "目标检测结果", "只覆盖96类IP102害虫\n可输出多框或无框\n框坐标还原到原图", 5.8, 3.55, 2.75, 1.8, fill=BLUE_LIGHT, line=BLUE)
    add_arrow(slide, 8.55, 2.15, 9.1, 2.15, GREEN)
    add_arrow(slide, 8.55, 4.45, 9.1, 4.45, BLUE)
    add_box(slide, "页面展示", "作物/四大类/细类\nTop-5", 9.05, 1.4, 1.9, 1.5, fill=WHITE, line=GREEN)
    add_box(slide, "页面叠框", "类别名/分数/位置", 9.05, 3.8, 1.9, 1.35, fill=WHITE, line=BLUE)
    add_box(slide, "明确限制", "病害与缺陷没有IP102框标注\n因此只能分类\n不能承诺定位", 11.05, 2.25, 1.7, 2.55, fill=RED_LIGHT, line=RED, title_size=15, subtitle_size=9.5)
    add_text(slide, "最终发布目标：一个 checkpoint，不再同时部署旧分类模型与单独检测模型", 1.2, 6.2, 10.9, 0.42, size=16, color=RED, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, 6)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)


if __name__ == "__main__":
    build()
